from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from tempfile import TemporaryDirectory
from typing import Iterable


@dataclass
class TextChunk:
    content: str
    page_number: int | None = None


@dataclass
class VisualUnit:
    kind: str
    path: Path
    page_number: int | None
    ocr_text: str


@dataclass
class ExtractedDocument:
    chunks: list[TextChunk]
    visuals: list[VisualUnit]
    temp_dir: TemporaryDirectory[str]


def extract_document(path: Path, filename: str, max_pdf_pages: int, render_pdf_pages: int) -> ExtractedDocument:
    temp_dir = TemporaryDirectory(prefix="voyager-worker-")
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        chunks, visuals = extract_pdf(path, Path(temp_dir.name), max_pdf_pages, render_pdf_pages)
    elif suffix in {".docx"}:
        chunks, visuals = chunk_plain_text(extract_docx(path)), []
    elif suffix in {".xlsx", ".xlsm"}:
        chunks, visuals = chunk_plain_text(extract_xlsx(path)), []
    elif suffix in {".pptx"}:
        chunks, visuals = chunk_plain_text(extract_pptx(path)), []
    elif suffix in {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff"}:
        ocr_text = ocr_image(path)
        chunks, visuals = chunk_plain_text(ocr_text or filename), [VisualUnit("image", path, None, ocr_text)]
    else:
        chunks, visuals = chunk_plain_text(path.read_text("utf-8", errors="ignore")), []

    if not chunks and visuals:
        chunks = [TextChunk((visuals[0].ocr_text or filename)[:4000], visuals[0].page_number)]
    return ExtractedDocument(chunks=chunks, visuals=visuals, temp_dir=temp_dir)


def extract_pdf(path: Path, temp_dir: Path, max_pages: int, render_pages: int) -> tuple[list[TextChunk], list[VisualUnit]]:
    import fitz

    document = fitz.open(path)
    chunks: list[TextChunk] = []
    visuals: list[VisualUnit] = []
    page_count = min(len(document), max_pages)
    for index in range(page_count):
        page = document[index]
        page_number = index + 1
        text = page.get_text("text").strip()
        if text:
            chunks.extend(chunk_plain_text(text, page_number))
        if index < render_pages:
            pixmap = page.get_pixmap(matrix=fitz.Matrix(1.4, 1.4), alpha=False)
            image_path = temp_dir / f"page-{page_number}.png"
            pixmap.save(image_path)
            ocr_text = "" if len(text) > 80 else ocr_image(image_path)
            if ocr_text:
                chunks.extend(chunk_plain_text(ocr_text, page_number))
            visuals.append(VisualUnit("page", image_path, page_number, ocr_text))
    return chunks, visuals


def extract_docx(path: Path) -> str:
    from docx import Document

    document = Document(str(path))
    return "\n".join(paragraph.text for paragraph in document.paragraphs if paragraph.text.strip())


def extract_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(path, data_only=True, read_only=True)
    try:
        lines: list[str] = []
        for sheet in workbook.worksheets:
            lines.append(f"# {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                values = [str(value) for value in row if value is not None and str(value).strip()]
                if values:
                    lines.append(" | ".join(values))
        return "\n".join(lines)
    finally:
        workbook.close()


def extract_pptx(path: Path) -> str:
    from pptx import Presentation

    presentation = Presentation(str(path))
    lines: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        lines.append(f"# Slide {index}")
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                lines.append(text.strip())
    return "\n".join(lines)


def ocr_image(path: Path) -> str:
    try:
        from rapidocr_onnxruntime import RapidOCR

        engine = RapidOCR()
        result, _ = engine(str(path))
        if not result:
            return ""
        return "\n".join(item[1] for item in result if len(item) > 1 and item[1])
    except Exception:
        return ""


def chunk_plain_text(text: str, page_number: int | None = None, max_chars: int = 1200, overlap: int = 160) -> list[TextChunk]:
    text = text.replace("\x00", "")
    cleaned = "\n".join(line.strip() for line in text.splitlines() if line.strip())
    if not cleaned:
        return []
    chunks: list[TextChunk] = []
    start = 0
    while start < len(cleaned):
        end = min(start + max_chars, len(cleaned))
        chunks.append(TextChunk(cleaned[start:end], page_number))
        if end == len(cleaned):
            break
        start = max(0, end - overlap)
    return chunks
